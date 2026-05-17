"""
KidGuard PowerPoint Presentation Builder
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1A, 0x23, 0x7E)   # deep navy
BLUE       = RGBColor(0x29, 0x79, 0xFF)   # bright blue
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)   # very light blue
RED        = RGBColor(0xD3, 0x2F, 0x2F)   # alert red
GREEN      = RGBColor(0x2E, 0x7D, 0x32)   # success green
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x37, 0x47, 0x4F)
MID_GRAY   = RGBColor(0x78, 0x90, 0x9C)
AMBER      = RGBColor(0xF5, 0x7F, 0x17)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper functions ────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, fill_alpha=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_multiline(slide, lines, x, y, w, h, font_size=16, color=DARK_GRAY,
                  bold_first=False, line_spacing=None):
    """lines = list of strings. Prefix line with '•' to auto-bullet."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox

def navy_header(slide, title, subtitle=None):
    """Full-width navy header bar at top."""
    add_rect(slide, 0, 0, 13.33, 1.3, fill_color=NAVY)
    add_text(slide, title, 0.4, 0.08, 12, 0.7, font_size=32, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.75, 12, 0.45, font_size=15,
                 color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

def slide_footer(slide, page_num, total=10):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill_color=NAVY)
    add_text(slide, "KidGuard  |  Child Safety Intelligence Platform",
             0.3, 7.22, 10, 0.25, font_size=9, color=LIGHT_BLUE)
    add_text(slide, f"{page_num} / {total}", 12.5, 7.22, 0.7, 0.25,
             font_size=9, color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, header_text, body_lines,
         header_color=NAVY, body_font=14):
    add_rect(slide, x, y, w, 0.45, fill_color=header_color)
    add_text(slide, header_text, x+0.1, y+0.06, w-0.2, 0.35,
             font_size=14, bold=True, color=WHITE)
    add_rect(slide, x, y+0.45, w, h-0.45, fill_color=WHITE,
             line_color=MID_GRAY, line_width=0.5)
    body = "\n".join(body_lines)
    add_text(slide, body, x+0.15, y+0.52, w-0.3, h-0.6,
             font_size=body_font, color=DARK_GRAY, align=PP_ALIGN.LEFT, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Hero
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# Full navy background
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=NAVY)

# Diagonal accent strip
accent = sl.shapes.add_shape(1, Inches(7.8), Inches(0), Inches(5.53), Inches(7.5))
accent.fill.solid()
accent.fill.fore_color.rgb = BLUE
accent.line.fill.background()

# Overlapping dark panel
add_rect(sl, 8.2, 0, 5.13, 7.5, fill_color=RGBColor(0x0D, 0x18, 0x5C))

# Logo / brand mark placeholder
add_rect(sl, 8.6, 1.2, 4.2, 4.2, fill_color=RGBColor(0x1E, 0x2B, 0x8A))
add_text(sl, "🛡", 9.6, 1.7, 2.2, 2.2, font_size=72, color=BLUE, align=PP_ALIGN.CENTER)
add_text(sl, "KidGuard", 8.6, 3.9, 4.2, 0.7, font_size=22, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Child Safety Intelligence", 8.6, 4.5, 4.2, 0.4, font_size=13,
         color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Title text
add_text(sl, "KidGuard", 0.5, 1.4, 7.5, 1.1, font_size=54, bold=True, color=WHITE)
add_text(sl, "AI-Powered Child Safety\nfor the Modern Internet", 0.5, 2.6, 7.2, 1.4,
         font_size=26, color=LIGHT_BLUE)
add_text(sl,
    "Real-time grooming detection on Roblox\n"
    "— built for parents, powered by AI.",
    0.5, 4.2, 7.2, 0.9, font_size=17, color=WHITE, italic=True)

# Bottom tagline bar
add_rect(sl, 0, 6.8, 8.1, 0.45, fill_color=BLUE)
add_text(sl, "Protecting children.  One conversation at a time.",
         0.4, 6.86, 7.5, 0.35, font_size=14, bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=RGBColor(0xFA, 0xFA, 0xFA))
navy_header(sl, "The Problem", "Online child grooming is growing — and most parents have no idea it's happening")
slide_footer(sl, 2)

stats = [
    ("500 M+", "Roblox monthly\nactive users"),
    ("67%", "are under\n16 years old"),
    ("1 in 7", "children receive\nonline sexual\nsolicitation"),
    ("93%", "of grooming\nstarts with\nnormal conversation"),
]
for i, (num, label) in enumerate(stats):
    cx = 0.4 + i * 3.2
    add_rect(sl, cx, 1.5, 2.9, 2.2, fill_color=NAVY)
    add_text(sl, num, cx, 1.58, 2.9, 1.1, font_size=38, bold=True,
             color=BLUE, align=PP_ALIGN.CENTER)
    add_text(sl, label, cx, 2.65, 2.9, 0.9, font_size=14,
             color=WHITE, align=PP_ALIGN.CENTER)

add_rect(sl, 0.4, 3.95, 12.5, 0.04, fill_color=BLUE)

add_multiline(sl, [
    "⚠  Children are targeted through legitimate gaming platforms — Roblox, Minecraft, Fortnite — where parents assume they are safe.",
    "⚠  Grooming is gradual. It begins with friendship, trust-building, and compliments before escalating.",
    "⚠  Most parents only discover grooming after serious harm has already occurred.",
    "⚠  Existing tools block content — they don't understand conversations.",
], 0.5, 4.1, 12.3, 2.9, font_size=15, color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Existing Solutions & Their Weaknesses
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=RGBColor(0xFA, 0xFA, 0xFA))
navy_header(sl, "What Else Exists — and Why It Falls Short",
            "Current solutions are either too broad, too slow, or too blind to context")
slide_footer(sl, 3)

competitors = [
    ("Net Nanny", RED,   ["✗  Full device agent install", "✗  Blocks websites — can't read chat", "✗  No AI understanding of context", "✗  Expensive monthly subscription", "✗  Not gaming-platform aware"]),
    ("Bark",      AMBER, ["✗  Monitors after the fact (not real-time)", "✗  Email/social focused — misses Roblox", "✗  No conversation context captured", "✗  No SMS instant alert", "✗  High false-positive rate"]),
    ("Qustodio",  AMBER, ["✗  Keyword-based filtering only", "✗  Easily bypassed with synonyms", "✗  No AI grooming pattern detection", "✗  Network-level only — misses encryption", "✗  Generic parental controls"]),
    ("Roblox\nChat Filter", MID_GRAY, ["✗  Only filters profanity, not grooming", "✗  No parent notification system", "✗  No conversation history", "✗  Bypass patterns widely known", "✗  No risk scoring or AI analysis"]),
]
for i, (name, color, bullets) in enumerate(competitors):
    cx = 0.3 + i * 3.25
    add_rect(sl, cx, 1.45, 3.0, 0.5, fill_color=color)
    add_text(sl, name, cx+0.1, 1.5, 2.8, 0.4, font_size=15, bold=True, color=WHITE)
    add_rect(sl, cx, 1.95, 3.0, 4.7, fill_color=WHITE,
             line_color=color, line_width=1.2)
    body = "\n".join(bullets)
    add_text(sl, body, cx+0.12, 2.05, 2.78, 4.5,
             font_size=12.5, color=DARK_GRAY, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Introducing KidGuard
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=NAVY)
add_rect(sl, 0, 0, 13.33, 1.3, fill_color=RGBColor(0x0D, 0x18, 0x5C))
add_text(sl, "Introducing KidGuard", 0.4, 0.1, 12, 0.7, font_size=32,
         bold=True, color=WHITE)
add_text(sl, "A lightweight, AI-native solution built specifically for gaming platforms",
         0.4, 0.78, 12, 0.4, font_size=15, color=LIGHT_BLUE)
slide_footer(sl, 4)

pillars = [
    ("🎯  Platform-Specific", "Built exclusively for Roblox — understands the context, the language, and the risk patterns unique to gaming chat."),
    ("🤖  AI-Native Detection", "Uses Claude AI (Anthropic) to understand conversation intent — not just keywords. Detects manipulation, trust-building, and bypass attempts."),
    ("⚡  Real-Time Alerts", "Parents receive an SMS the moment a risk is detected — not hours later when the damage is done."),
    ("🪶  Zero Installation", "A simple Chrome extension. No device agents, no VPNs, no router changes. Works in minutes."),
    ("📋  Evidence-Ready", "Captures 10-message conversation context at the moment of detection — ready for authorities if needed."),
    ("🔒  Privacy First", "Monitors only Roblox chat. No screenshots, no full device surveillance."),
]
for i, (title, body) in enumerate(pillars):
    row = i // 3
    col = i  % 3
    cx = 0.3 + col * 4.34
    cy = 1.5 + row * 2.7
    add_rect(sl, cx, cy, 4.1, 2.45, fill_color=RGBColor(0x1E, 0x2B, 0x8A))
    add_text(sl, title, cx+0.15, cy+0.12, 3.8, 0.45, font_size=14,
             bold=True, color=BLUE)
    add_text(sl, body, cx+0.15, cy+0.55, 3.8, 1.75, font_size=12.5,
             color=WHITE, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — How KidGuard is Better
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=RGBColor(0xFA, 0xFA, 0xFA))
navy_header(sl, "KidGuard vs The Competition",
            "Not just better — fundamentally different in approach")
slide_footer(sl, 5)

# Comparison table
headers = ["Feature", "Net Nanny", "Bark", "Roblox Filter", "KidGuard ✓"]
col_w   = [3.2, 2.1, 2.1, 2.1, 2.6]
col_x   = [0.3]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w + 0.08)

header_colors = [NAVY, MID_GRAY, MID_GRAY, MID_GRAY, GREEN]
for j, (h, cx, cw, hc) in enumerate(zip(headers, col_x, col_w, header_colors)):
    add_rect(sl, cx, 1.45, cw, 0.45, fill_color=hc)
    add_text(sl, h, cx+0.05, 1.5, cw-0.1, 0.35, font_size=13,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Real-time detection",          "✗", "✗", "✗", "✓"),
    ("AI conversation understanding","✗", "Partial", "✗", "✓"),
    ("Grooming bypass detection",    "✗", "✗", "✗", "✓"),
    ("SMS instant parent alert",     "✓", "✓", "✗", "✓"),
    ("Conversation evidence capture","✗", "✗", "✗", "✓"),
    ("No device install needed",     "✗", "✗", "✓", "✓"),
    ("Roblox-specific context",      "✗", "✗", "Partial", "✓"),
    ("Free / low cost",              "✗", "✗", "✓", "✓"),
]
row_colors = [RGBColor(0xFF,0xFF,0xFF), RGBColor(0xF1,0xF8,0xFE)]
for r, (feat, *vals) in enumerate(rows):
    ry = 1.93 + r * 0.56
    bg = row_colors[r % 2]
    for j, (val, cx, cw) in enumerate(zip([feat]+vals, col_x, col_w)):
        add_rect(sl, cx, ry, cw, 0.52, fill_color=bg,
                 line_color=RGBColor(0xDD,0xDD,0xDD), line_width=0.3)
        fc = GREEN if val == "✓" else (RED if val == "✗" else AMBER)
        fc = DARK_GRAY if j == 0 else fc
        add_text(sl, val, cx+0.05, ry+0.12, cw-0.1, 0.3,
                 font_size=13, color=fc, align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT,
                 bold=(j==4))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Architecture
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=RGBColor(0xF5, 0xF7, 0xFF))
navy_header(sl, "System Architecture", "Lightweight, modular, and cloud-ready")
slide_footer(sl, 6)

# Flow boxes
boxes = [
    (0.25, 2.5, 2.3, 3.5,  NAVY,  "Chrome\nExtension",
     "• MutationObserver\n• Pattern scorer\n• 10-msg context\n• Badge alerts"),
    (3.15, 2.5, 2.5, 3.5,  BLUE,  "FastAPI\nBackend",
     "• REST API\n• Incident routing\n• Risk upgrade logic\n• Auth / JWT"),
    (6.15, 1.5, 2.5, 2.0,  RGBColor(0x6A,0x1B,0x9A), "Claude Haiku AI",
     "• Grooming classifier\n• 17/17 test accuracy\n• Context-aware"),
    (6.15, 4.0, 2.5, 2.0,  RGBColor(0x00,0x69,0x5C), "MongoDB",
     "• Incident storage\n• Conversation context\n• Parent accounts"),
    (9.2,  2.0, 2.3, 1.6,  RGBColor(0xD8,0x43,0x15), "Twilio SMS",
     "• Instant parent alert\n• A2P 10DLC"),
    (9.2,  4.0, 2.3, 1.6,  RGBColor(0x18,0x65,0x38), "React Dashboard",
     "• Incident review\n• Evidence viewer"),
    (9.2,  6.0, 2.3, 1.0,  MID_GRAY, "Future: Cloud",
     "• AWS / GCP deploy"),
]

for (bx, by, bw, bh, bc, btitle, bbody) in boxes:
    add_rect(sl, bx, by, bw, 0.5, fill_color=bc)
    add_text(sl, btitle, bx+0.08, by+0.06, bw-0.16, 0.4,
             font_size=13, bold=True, color=WHITE)
    add_rect(sl, bx, by+0.5, bw, bh-0.5, fill_color=WHITE,
             line_color=bc, line_width=1.2)
    add_text(sl, bbody, bx+0.12, by+0.58, bw-0.24, bh-0.62,
             font_size=11.5, color=DARK_GRAY, wrap=True)

# Arrows (simple lines)
def arrow(slide, x1, y1, x2, y2):
    from pptx.util import Inches
    connector = slide.shapes.add_connector(1,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = MID_GRAY
    connector.line.width = Pt(1.8)

arrow(sl, 2.55, 4.25, 3.15, 4.25)   # extension → backend
arrow(sl, 5.65, 3.5, 6.15, 2.5)     # backend → AI
arrow(sl, 5.65, 4.5, 6.15, 5.0)     # backend → mongo
arrow(sl, 8.65, 2.5, 9.2, 2.8)      # AI → twilio
arrow(sl, 8.65, 5.0, 9.2, 4.8)      # mongo → dashboard

# Roblox user label
add_rect(sl, 0.25, 1.6, 2.3, 0.65, fill_color=RGBColor(0xE3,0xF2,0xFD))
add_text(sl, "👦  Child on Roblox", 0.3, 1.66, 2.2, 0.45,
         font_size=12, color=NAVY, align=PP_ALIGN.CENTER)
arrow(sl, 1.4, 2.25, 1.4, 2.5)      # roblox → extension

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — The AI Detection Layer
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=NAVY)
add_rect(sl, 0, 0, 13.33, 1.3, fill_color=RGBColor(0x0D, 0x18, 0x5C))
add_text(sl, "The AI Detection Layer", 0.4, 0.1, 12, 0.7,
         font_size=32, bold=True, color=WHITE)
add_text(sl, "Claude Haiku — purpose-built for Roblox child safety classification",
         0.4, 0.78, 12, 0.4, font_size=15, color=LIGHT_BLUE)
slide_footer(sl, 7)

# Accuracy callout
add_rect(sl, 0.3, 1.5, 3.8, 2.5, fill_color=RGBColor(0x1E, 0x2B, 0x8A))
add_text(sl, "17 / 17", 0.3, 1.55, 3.8, 1.3, font_size=52,
         bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(sl, "Test accuracy\non real grooming\nscenarios",
         0.3, 2.85, 3.8, 0.9, font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(sl, 0.3, 4.2, 3.8, 1.6, fill_color=RGBColor(0x1E, 0x2B, 0x8A))
add_text(sl, "Detects bypass\nattempts", 0.5, 4.28, 3.6, 0.45,
         font_size=14, bold=True, color=BLUE)
add_text(sl,
    '"h3y wanna b fr1ends?"\n"can we talk on discord instead?"\n"don\'t tell your parents about me"',
    0.5, 4.72, 3.6, 0.95, font_size=11.5, color=MID_GRAY, italic=True)

# What the AI does
add_rect(sl, 4.4, 1.5, 8.6, 5.7, fill_color=RGBColor(0x1A, 0x27, 0x72))

points = [
    ("Context-aware, not keyword-based",
     "Understands that 'are you alone?' means something very different in a grooming context than in a normal conversation."),
    ("Roblox-specific training context",
     "System prompt grounds every classification in the Roblox gaming environment — reducing false positives from normal gaming chat."),
    ("Three-tier risk output",
     "Returns GROOMING / SUSPICIOUS / SAFE with a confidence score (0-1) and a human-readable reason for parent review."),
    ("Risk upgrade logic",
     "AI confidence ≥ 0.9 → upgrades incident to HIGH risk.  Confidence ≥ 0.8 → MEDIUM.  Layered defence with rule-based scorer."),
    ("Evidence-quality reasoning",
     "Every detection includes a written reason — e.g. 'Asking child to move conversation off-platform is a classic isolation tactic.'"),
]
for i, (title, body) in enumerate(points):
    ty = 1.65 + i * 1.06
    add_text(sl, f"✦  {title}", 4.6, ty, 8.1, 0.35,
             font_size=14, bold=True, color=BLUE)
    add_text(sl, body, 4.6, ty+0.35, 8.1, 0.6,
             font_size=12, color=WHITE, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — What Makes It Amazing
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=RGBColor(0xFA, 0xFA, 0xFA))
navy_header(sl, "What Makes KidGuard Remarkable",
            "A combination of innovations that no existing product offers")
slide_footer(sl, 8)

wows = [
    (NAVY,  "🎯", "Zero-Install\nDeployment",
     "A Chrome extension. Parents install it in under 2 minutes — no IT skills, no router config, no VPN."),
    (BLUE,  "🤖", "AI That Actually\nUnderstands",
     "Most tools match keywords. KidGuard reads the room — detecting manipulation patterns, urgency, secrecy requests."),
    (GREEN, "⚡", "Sub-Second\nParent Alert",
     "From chat message to parent SMS in under 2 seconds. Real-time response during the critical window."),
    (RGBColor(0x6A,0x1B,0x9A), "📋", "Court-Ready\nEvidence",
     "10-message conversation context automatically captured — structured, timestamped, downloadable JSON evidence."),
    (RED,   "🕵️", "Bypass-Proof\nDetection",
     "Predators use leetspeak, coded language, and platform-switching tactics. KidGuard is trained to spot all of them."),
    (AMBER, "🔒", "Privacy by\nDesign",
     "Only monitors Roblox chat. No screenshots, no keystroke logging, no browsing history. Surgical and ethical."),
]
for i, (color, icon, title, body) in enumerate(wows):
    row = i // 3
    col = i  % 3
    cx = 0.3 + col * 4.34
    cy = 1.5 + row * 2.8
    add_rect(sl, cx, cy, 4.1, 2.55, fill_color=WHITE,
             line_color=color, line_width=2)
    add_rect(sl, cx, cy, 0.55, 2.55, fill_color=color)
    add_text(sl, icon, cx+0.02, cy+0.8, 0.55, 0.7, font_size=20,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, cx+0.7, cy+0.18, 3.25, 0.65,
             font_size=15, bold=True, color=color)
    add_text(sl, body, cx+0.7, cy+0.85, 3.25, 1.5,
             font_size=12.5, color=DARK_GRAY, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — How It Works (User Journey)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=RGBColor(0xF0, 0xF4, 0xFF))
navy_header(sl, "How It Works", "From Roblox chat to parent awareness in real time")
slide_footer(sl, 9)

steps = [
    ("1", NAVY,  "Child Opens\nRoblox",
     "Extension activates automatically on roblox.com. No action needed from the child."),
    ("2", BLUE,  "Chat is\nMonitored",
     "MutationObserver watches for new messages in real time. Existing chat is skipped — no false replay."),
    ("3", RGBColor(0x6A,0x1B,0x9A), "Risk Scored\nLocally",
     "Built-in pattern engine scores the message (0-10). Suspicious messages are escalated to AI."),
    ("4", RGBColor(0xD8,0x43,0x15), "AI Classifies\nthe Threat",
     "Claude Haiku reviews the message and conversation context. Returns label + confidence + reason."),
    ("5", RED,   "Incident\nLogged",
     "Incident stored with 10-message context. Risk level upgraded if AI confidence is high."),
    ("6", GREEN, "Parent\nAlerted",
     "SMS sent instantly via Twilio. Parent opens dashboard to review the full conversation evidence."),
]

for i, (num, color, title, body) in enumerate(steps):
    cx = 0.3 + i * 2.14
    # connector arrow
    if i < 5:
        add_rect(sl, cx+2.0, 2.45, 0.2, 0.12, fill_color=MID_GRAY)
    add_rect(sl, cx, 1.45, 1.95, 0.6, fill_color=color)
    add_text(sl, f"STEP {num}", cx+0.07, 1.52, 1.82, 0.45,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, cx, 2.05, 1.95, 4.55, fill_color=WHITE,
             line_color=color, line_width=1.5)
    add_text(sl, title, cx+0.1, 2.15, 1.78, 0.7,
             font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(sl, body, cx+0.1, 2.9, 1.78, 3.5,
             font_size=12, color=DARK_GRAY, wrap=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Vision & Call to Action
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=NAVY)
add_rect(sl, 0, 0, 13.33, 0.08, fill_color=BLUE)
add_rect(sl, 0, 7.42, 13.33, 0.08, fill_color=BLUE)

add_text(sl, "The Vision", 0.5, 0.4, 12.5, 0.75,
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "KidGuard is not just a Chrome extension — it is a platform.",
         0.5, 1.25, 12.5, 0.5, font_size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER,
         italic=True)

roadmap = [
    ("Now",       BLUE,  "Roblox grooming detection\nSMS alerts + evidence capture\nAI-powered classification"),
    ("Next",      AMBER, "Minecraft & Discord support\nMobile companion app\nParent dashboard mobile view"),
    ("Future",    GREEN, "Multi-platform AI safety layer\nSchool district licensing\nLaw enforcement integration"),
]
for i, (phase, color, items) in enumerate(roadmap):
    cx = 0.6 + i * 4.2
    add_rect(sl, cx, 2.0, 3.8, 0.5, fill_color=color)
    add_text(sl, phase, cx, 2.05, 3.8, 0.4,
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, cx, 2.5, 3.8, 2.4, fill_color=RGBColor(0x1E, 0x2B, 0x8A))
    add_text(sl, items, cx+0.15, 2.6, 3.5, 2.2,
             font_size=13, color=WHITE, wrap=True)

add_rect(sl, 0.5, 5.1, 12.3, 0.04, fill_color=BLUE)

add_text(sl,
    '"Every child deserves to play freely online. Every parent deserves to know they are safe."',
    0.5, 5.25, 12.5, 0.75, font_size=17, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

add_rect(sl, 3.5, 6.1, 6.3, 0.7, fill_color=BLUE)
add_text(sl, "🛡  KidGuard — Built for Today. Ready for Tomorrow.",
         3.5, 6.15, 6.3, 0.55, font_size=15, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
output_path = r"c:\Users\Nathan\Documents\Personal Learning\app\KidGuard_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
